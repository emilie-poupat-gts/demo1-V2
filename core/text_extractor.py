import pptx
import openpyxl
from PyPDF2 import PdfReader
from docx import Document

def extract_text(uploaded_file):
    ext = uploaded_file.name.lower().split(".")[-1]

    if ext == "pdf":
        reader = PdfReader(uploaded_file)
        return "\n".join([p.extract_text() or "" for p in reader.pages])

    if ext == "docx":
        doc = Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs])

    if ext == "pptx":
        prs = pptx.Presentation(uploaded_file)
        return "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    if ext in ["xlsx", "xls"]:
        wb = openpyxl.load_workbook(uploaded_file, data_only=True)
        texte = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                texte += " ".join([str(c) for c in row if c]) + "\n"
        return texte
    
    if ext=="txt":
        return uploaded_file.read().decode("utf-8", errors="ignore")

    else:
        raise ValueError(f"Unsupported file type: {suffix}")
